"""
PowerPoint Parser for PPTX files
Extracts text, tables, bullet points, and images from slides
Adapted from RAGFlow's ppt_parser.py

Images are extracted and returned in the same format as DoclingParser,
allowing process_document.py to describe them with FigureParser/GPT-4o Vision.
"""

import io
import logging
from io import BytesIO
from typing import Dict, Any, List

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    Presentation = None


class PPTParser:
    """Parser for PowerPoint files using python-pptx"""

    SUPPORTED_FORMATS = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }

    def can_parse(self, content_type: str) -> bool:
        """Check if this parser can handle the content type"""
        if not content_type:
            return False
        return content_type in self.SUPPORTED_FORMATS

    def _get_bulleted_text(self, paragraph) -> str:
        """Extract text with bullet point formatting preserved"""
        try:
            is_bulleted = (
                bool(paragraph._p.xpath("./a:pPr/a:buChar")) or
                bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or
                bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
            )
            if is_bulleted:
                indent = "  " * paragraph.level
                return f"{indent}- {paragraph.text}"
            return paragraph.text
        except Exception:
            return paragraph.text if hasattr(paragraph, 'text') else ""

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape, handling various shape types"""
        try:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self._get_bulleted_text(paragraph))
                return "\n".join(texts)

            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                if hasattr(shape, 'text'):
                    return shape.text.strip()
                return ""

            # Handle table (shape_type 19)
            if shape_type == 19:
                return self._extract_table(shape)

            # Handle group shape (shape_type 6)
            if shape_type == 6:
                texts = []
                for p in sorted(shape.shapes, key=lambda x: (
                    x.top // 10 if x.top else 0,
                    x.left if x.left else 0
                )):
                    t = self._extract_shape_text(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)

            return ""

        except Exception as e:
            logger.debug(f"Error extracting shape text: {e}")
            return ""

    def _extract_table(self, shape) -> str:
        """Extract table content as structured text"""
        try:
            tb = shape.table
            rows = []

            # Get header row
            if len(tb.rows) > 0:
                headers = [tb.cell(0, j).text for j in range(len(tb.columns))]

                # Data rows with header context
                for i in range(1, len(tb.rows)):
                    row_data = []
                    for j in range(len(tb.columns)):
                        cell_text = tb.cell(i, j).text
                        if cell_text.strip():
                            if headers[j].strip():
                                row_data.append(f"{headers[j]}: {cell_text}")
                            else:
                                row_data.append(cell_text)
                    if row_data:
                        rows.append("; ".join(row_data))

            return "\n".join(rows)
        except Exception as e:
            logger.debug(f"Error extracting table: {e}")
            return ""

    def _extract_images_from_slide(
        self, slide, slide_number: int, image_index_start: int
    ) -> List[Dict[str, Any]]:
        """
        Extract all images from a slide.

        Args:
            slide: python-pptx slide object
            slide_number: 1-indexed slide number (used as page)
            image_index_start: Starting index for image numbering

        Returns:
            List of image dicts matching DoclingParser format:
            - data: Raw image bytes
            - page: Slide number
            - index: Global image index
            - format: Image format (png, jpeg, etc.)
            - filename: Generated filename
        """
        images = []
        img_idx = image_index_start

        for shape in slide.shapes:
            try:
                # Check if shape is a picture (shape_type 13 = MSO_SHAPE_TYPE.PICTURE)
                try:
                    shape_type = shape.shape_type
                except NotImplementedError:
                    continue

                # shape_type 13 is PICTURE in python-pptx
                if shape_type == 13 and hasattr(shape, "image"):
                    image = shape.image

                    # Get image bytes
                    image_bytes = image.blob

                    # Determine format from content type
                    content_type = image.content_type or "image/png"
                    if "jpeg" in content_type or "jpg" in content_type:
                        fmt = "jpeg"
                    elif "png" in content_type:
                        fmt = "png"
                    elif "gif" in content_type:
                        fmt = "gif"
                    elif "webp" in content_type:
                        fmt = "webp"
                    else:
                        fmt = "png"  # Default

                    images.append({
                        "data": image_bytes,
                        "page": slide_number,
                        "index": img_idx,
                        "format": fmt,
                        "filename": f"slide_{slide_number}_image_{img_idx}.{fmt}",
                    })
                    img_idx += 1
                    logger.debug(
                        f"Extracted image {img_idx} from slide {slide_number}"
                    )

            except Exception as e:
                logger.debug(f"Error extracting image from slide {slide_number}: {e}")
                continue

        return images

    def _parse_slide(self, slide, slide_number: int) -> Dict[str, Any]:
        """Parse a single slide and extract content"""
        texts = []

        # Sort shapes by position (top to bottom, left to right)
        sorted_shapes = sorted(
            slide.shapes,
            key=lambda x: (
                (x.top if x.top is not None else 0) // 10,
                x.left if x.left is not None else 0
            )
        )

        for shape in sorted_shapes:
            try:
                txt = self._extract_shape_text(shape)
                if txt:
                    texts.append(txt)
            except Exception as e:
                logger.debug(f"Error processing shape on slide {slide_number}: {e}")

        return {
            "slide_number": slide_number,
            "content": "\n".join(texts),
        }

    async def parse(self, file_path: str) -> Dict[str, Any]:
        """
        Parse PowerPoint file and extract content from all slides.

        Args:
            file_path: Path to PPTX file

        Returns:
            Dict with:
                - content: Combined text from all slides
                - metadata: Slide count and per-slide info
                - page_count: Number of slides
                - images: List of extracted images (same format as DoclingParser)
                    Each image dict contains:
                    - data: Raw image bytes
                    - page: Slide number
                    - index: Global image index
                    - format: Image format (png, jpeg, etc.)
                    - filename: Generated filename
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PowerPoint parsing. "
                "Install with: pip install python-pptx"
            )

        # Open presentation
        if isinstance(file_path, bytes):
            ppt = Presentation(BytesIO(file_path))
        else:
            ppt = Presentation(file_path)

        slides_data = []
        content_parts = []
        all_images = []
        image_index = 0

        for i, slide in enumerate(ppt.slides):
            slide_number = i + 1

            # Extract text content
            slide_data = self._parse_slide(slide, slide_number)
            slides_data.append(slide_data)

            if slide_data["content"]:
                content_parts.append(f"## Slide {slide_number}\n")
                content_parts.append(slide_data["content"])
                content_parts.append("")  # Empty line between slides

            # Extract images from this slide
            slide_images = self._extract_images_from_slide(
                slide, slide_number, image_index
            )
            all_images.extend(slide_images)
            image_index += len(slide_images)

        content = "\n".join(content_parts)

        # Log image extraction results
        if all_images:
            logger.info(
                f"Extracted {len(all_images)} images from {len(ppt.slides)} slides"
            )

        metadata = {
            "slide_count": len(ppt.slides),
            "image_count": len(all_images),
            "slides": [
                {
                    "slide_number": s["slide_number"],
                    "has_content": bool(s["content"]),
                    "content_length": len(s["content"]),
                }
                for s in slides_data
            ],
        }

        return {
            "content": content,
            "metadata": metadata,
            "page_count": len(ppt.slides),
            "images": all_images,
        }
